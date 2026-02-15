# ==========================================
# 라이브러리 불러오기 (import)
# ==========================================
import streamlit as st
from google import genai
import json
import os
from datetime import datetime, timedelta
import docx
from pypdf import PdfReader
import random
from pptx import Presentation as PptxPresentation
import pandas as pd
from io import BytesIO
from docx import Document as DocxDocument
from docx.shared import Cm, Pt, RGBColor
from docx.enum.text import WD_COLOR_INDEX, WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_ROW_HEIGHT_RULE, WD_CELL_VERTICAL_ALIGNMENT 
from docx.oxml.ns import qn 
from docx.oxml import OxmlElement
import re

# ==========================================
# 1. 프로그램 기본 설정
# ==========================================
# [주의] 배포 시에는 st.secrets를 사용하세요.
GOOGLE_API_KEY = st.secrets["GOOGLE_API_KEY"]
client = genai.Client(api_key=GOOGLE_API_KEY)
MODEL = 'gemini-2.5-flash'
DB_FILE = "medical_flashcards.json"

st.set_page_config(page_title="MEDI-Quiz", page_icon="🩺", layout="wide")

# ==========================================
# CSS 스타일 설정
# ==========================================
st.markdown("""
<style>
    /* 1. 전체 폰트 및 기본 스타일 */
    .question-box {
        background-color: #f8f9fa; padding: 25px; border-radius: 12px; 
        border: 1px solid #e9ecef; margin-bottom: 25px; font-size: 1.1rem; line-height: 1.6;
    }
    .options-box {
        background-color: #f8f9fa; padding: 20px; border-radius: 12px; 
        border: 1px solid #e9ecef; margin-bottom: 25px;
    }
    .option-item {
        display: flex; align-items: center; padding: 12px 15px; 
        margin-bottom: 10px; border-radius: 8px; transition: background-color 0.2s;
    }
    .option-item:hover { background-color: #e9ecef; }
    .option-number {
        font-size: 1.1rem; font-weight: bold; margin-right: 15px; min-width: 30px;
    }
    .eliminated { text-decoration: line-through; color: #adb5bd; }
    
    /* 2. 하이라이트 스타일 */
    .hl-yellow { background-color: #fff3bf; padding: 2px 4px; border-radius: 3px; }
    .hl-blue { color: #1971c2; font-weight: bold; }
    .hl-gray { color: #adb5bd; }

    /* 3. 탭 스타일 (4등분, 가운데 정렬) */
    [data-testid="stTabs"] [role="tablist"] { display: flex !important; width: 100% !important; }
    [data-testid="stTabs"] button[role="tab"] { flex: 1 1 25% !important; justify-content: center !important; }
    [data-testid="stTabs"] button[role="tab"] p { font-size: 1.3rem !important; text-align: center !important; }

    /* 4. 파일 업로더 디자인 커스터마이징 */
    [data-testid="stFileUploader"] { margin-top: 20px; }
    [data-testid="stFileUploaderDropzone"] {
        background-color: #fff8f5;
        border: 2px dashed #FF6B35 !important;
        border-radius: 12px;
        padding: 40px 20px;
        min-height: 500px;
        display: flex;
        flex-direction: column;
        justify-content: center;
        align-items: center;
        text-align: center;
        transition: background-color 0.3s;
    }
    [data-testid="stFileUploaderDropzone"]:hover { background-color: #ffe8cc; }
    [data-testid="stColumn"]:first-child [data-testid="stFileUploaderDropzone"]::before {
        content: "📝"; font-size: 5rem; margin-bottom: 10px; display: block;
    }
    [data-testid="stColumn"]:first-child [data-testid="stFileUploaderDropzone"]::after {
        content: "정리본 / 강의자료를 업로드하세요";
        white-space: pre-wrap; font-size: 1.2rem; color: #495057; margin-top: 15px; font-weight: 600; line-height: 1.6;
    }
    [data-testid="stColumn"]:last-child [data-testid="stFileUploaderDropzone"]::before {
        content: "🏆"; font-size: 5rem; margin-bottom: 10px; display: block;
    }
    [data-testid="stColumn"]:last-child [data-testid="stFileUploaderDropzone"]::after {
        content: "족보 (선택사항)";
        white-space: pre-wrap; font-size: 1.2rem; color: #495057; margin-top: 15px; font-weight: 600; line-height: 1.6;
    }
    [data-testid="stFileUploaderDropzoneInstructions"], [data-testid="stFileUploaderDropzone"] small { display: none !important; }
    [data-testid="stFileUploaderDropzone"] button {
        background-color: #FF6B35; color: white; border: none; border-radius: 20px; padding: 10px 25px; font-weight: bold; order: 2; font-size: 1rem;
    }
    [data-testid="stFileUploaderDropzone"] button:hover { background-color: #e8590c; color: white; }
</style>
""", unsafe_allow_html=True)

# ==========================================
# 2. 백엔드 함수들
# ==========================================

def set_cell_background(cell, color_hex):
    cell_properties = cell._element.get_or_add_tcPr()
    shading_elm = OxmlElement('w:shd')
    shading_elm.set(qn('w:fill'), color_hex)
    cell_properties.append(shading_elm)

def set_font_style(run, font_name='맑은 고딕', font_size=9, is_bold=False):
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.bold = is_bold
    r = run._element
    rPr = r.get_or_add_rPr()
    fonts = OxmlElement('w:rFonts')
    fonts.set(qn('w:eastAsia'), font_name) 
    fonts.set(qn('w:ascii'), font_name)    
    fonts.set(qn('w:hAnsi'), font_name)    
    rPr.append(fonts)

def load_cards():
    if not os.path.exists(DB_FILE): return []
    with open(DB_FILE, "r", encoding="utf-8") as f:
        try:
            data = json.load(f)
            return [card for card in data if 'options' in card and isinstance(card['options'], list)]
        except: return []

def save_all_cards(cards):
    with open(DB_FILE, "w", encoding="utf-8") as f:
        json.dump(cards, f, ensure_ascii=False, indent=4)

def save_card_to_file(question, options, correct_index, explanation):
    cards = load_cards()
    cards.append({
        "question": question, "options": options, "correct_index": correct_index,
        "explanation": explanation, "next_review": datetime.now().strftime("%Y-%m-%d"), "interval": 1
    })
    save_all_cards(cards)

def delete_card(index):
    cards = load_cards()
    if 0 <= index < len(cards): del cards[index]; save_all_cards(cards)

def update_card_schedule(card_index, is_correct):
    cards = load_cards()
    if card_index < len(cards):
        card = cards[card_index]
        if is_correct:
            card['interval'] = card['interval'] * 2 + 1
            st.toast(f"🎉 정답! {card['interval']}일 뒤에 봅니다.")
        else:
            card['interval'] = 1
            st.toast("🥲 오답... 내일 다시 복습!")
        card['next_review'] = (datetime.now() + timedelta(days=card['interval'])).strftime("%Y-%m-%d")
        save_all_cards(cards)

def read_file(file):
    try:
        if file.name.endswith('.pdf'):
            reader = PdfReader(file)
            return "\n".join([page.extract_text() for page in reader.pages])
        elif file.name.endswith('.docx'):
            doc = docx.Document(file)
            return "\n".join([para.text for para in doc.paragraphs])
    except: return ""
    return ""

# ==========================================
# 3. 화면 구성
# ==========================================
st.markdown("<h1 style='text-align: center; color: #FF6B35; font-size: 3.2rem; margin-bottom: 2em;'>MEDI-Quiz</h1>", unsafe_allow_html=True)

if 'generated_quiz' not in st.session_state: st.session_state['generated_quiz'] = None
if 'show_explanation' not in st.session_state: st.session_state['show_explanation'] = False
if 'summary_data' not in st.session_state: st.session_state['summary_data'] = None

tab4, tab1, tab2, tab3 = st.tabs(["📋 정리본 형성", "📝 문제 생성", "🧠 실전 모의고사", "🗂️ 문제 관리"])

# ==========================================
# [탭 1] 문제 생성 (AI 쫄보 방지 및 5문제 강제 출제)
# ==========================================
with tab1:
    quiz_note_content = ""
    quiz_jokbo_content = ""
    col_q1, col_q2 = st.columns(2)

    with col_q1:
        quiz_note_file = st.file_uploader("정리본 업로드", type=['docx', 'pdf', 'pptx'], key="quiz_note_uploader", label_visibility="collapsed")
        if quiz_note_file:
            if quiz_note_file.name.endswith('.pptx'):
                try:
                    prs = PptxPresentation(quiz_note_file)
                    txt = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame: txt.append(shape.text_frame.text)
                    quiz_note_content = "\n".join(txt)
                except: pass
            else:
                quiz_note_content = read_file(quiz_note_file)
            if quiz_note_content:
                st.success(f"정리본 읽기 성공! ({len(quiz_note_content)}자)")


    with col_q2:
        quiz_jokbo_file = st.file_uploader("족보 업로드", type=['docx', 'pdf'], key="quiz_jokbo_uploader", label_visibility="collapsed")
        if quiz_jokbo_file:
            quiz_jokbo_content = read_file(quiz_jokbo_file)
            if quiz_jokbo_content:
                st.success(f"족보 읽기 성공! ({len(quiz_jokbo_content)}자)")

    st.divider()
    has_jokbo = bool(quiz_jokbo_content)

    if st.button("⚡ 5문제 출제하기", type="primary", use_container_width=True, disabled=not bool(quiz_note_content)):
        spinner_msg = "족보의 형식을 벤치마킹하여 정리본에서 5문제를 꽉 채워 출제 중입니다..." if has_jokbo else "정리본을 바탕으로 5문제를 만들고 있습니다..."
        with st.spinner(spinner_msg):
            try:
                if has_jokbo:
                    prompt = f"""
                    아래는 의대생이 공부한 정리본입니다. 이 학생이 정리본의 내용을 제대로 암기했는지 테스트하는 객관식 문제 5개를 만드세요.

                    [규칙]
                    - 정리본에 직접 나오는 질환명, 증상, 진단법, 치료법, 수치 등을 묻는 문제를 만드세요.
                    - "만약~했다면", "어떤 유형의 지식을~" 같은 메타 질문은 절대 만들지 마세요.
                    - 예시: "~의 1차 치료제는?", "~에서 나타나는 특징적 소견은?", "~의 진단 기준으로 옳은 것은?"
                    - [족보]의 문제 형식(문체, 보기 개수)만 참고하세요.

                    [정리본]
                    {quiz_note_content[:15000]}

                    [족보 - 형식 참고용]
                    {quiz_jokbo_content[:20000]}

                    JSON 배열로 5개 출력:
                    [{{"question": "질문", "options": ["보기1", "보기2", ...], "correct_index": 0, "explanation": "해설"}}]
                    """
                else:
                    prompt = f"""
                    아래는 의대생이 공부한 정리본입니다. 이 학생이 정리본의 내용을 제대로 암기했는지 테스트하는 5지선다형 객관식 문제 5개를 만드세요.

                    [규칙]
                    - 정리본에 직접 나오는 질환명, 증상, 진단법, 치료법, 수치 등을 묻는 문제를 만드세요.
                    - "만약~했다면", "어떤 유형의 지식을~" 같은 메타 질문은 절대 만들지 마세요.
                    - 예시: "~의 1차 치료제는?", "~에서 나타나는 특징적 소견은?", "~의 진단 기준으로 옳은 것은?"

                    [정리본]
                    {quiz_note_content[:15000]}

                    JSON 배열로 5개 출력:
                    [{{"question": "질문", "options": ["보기1", "보기2", "보기3", "보기4", "보기5"], "correct_index": 0, "explanation": "해설"}}]
                    """
                
                # 강제 JSON 출력 옵션 유지
                response = client.models.generate_content(
                    model=MODEL, 
                    contents=prompt,
                    config={"response_mime_type": "application/json"}
                )
                
                quizzes = json.loads(response.text)

                if isinstance(quizzes, list) and len(quizzes) > 0:
                    for quiz in quizzes:
                        save_card_to_file(quiz['question'], quiz['options'], quiz['correct_index'], quiz['explanation'])
                    st.success(f"✅ {len(quizzes)}개 문제가 생성되어 저장되었습니다! '실전 모의고사' 탭에서 확인하세요.")
                else: 
                    st.error("형식 오류: AI가 문제를 생성하지 못하고 빈 배열을 반환했습니다. 정리본 내용을 조금 더 추가해 보세요.")
            except json.JSONDecodeError:
                st.error("AI가 올바른 형식(JSON)으로 문제를 만들지 못했습니다. 다시 시도해 주세요.")
                with st.expander("AI 응답 원본 확인 (디버깅용)"):
                    st.write(response.text if 'response' in locals() else "응답 없음")
            except Exception as e: 
                st.error(f"오류: {e}")

# ==========================================
# [탭 2] 실전 모의고사
# ==========================================
with tab2:
    cards = load_cards()
    today = datetime.now().strftime("%Y-%m-%d")
    due_cards = [(i, c) for i, c in enumerate(cards) if c['next_review'] <= today]

    if not due_cards:
        st.info("🎉 오늘 풀 문제가 없습니다!")
    else:
        idx, card = due_cards[0]
        if 'current_quiz_idx' not in st.session_state or st.session_state.current_quiz_idx != idx:
            st.session_state.current_quiz_idx = idx
            st.session_state.selected_opt = None
            st.session_state.eliminated_opts = set()
            st.session_state.show_explanation = False

        st.write(f"남은 문제: **{len(due_cards)}개**")
        st.markdown(f"""<div class="question-box"><b>Q.</b> {card['question']}</div>""", unsafe_allow_html=True)
        st.write("---")

        circle_numbers = ["①", "②", "③", "④", "⑤"]
        st.markdown('<div class="options-box">', unsafe_allow_html=True)

        for i, opt_text in enumerate(card['options']):
            col_num, col_text, col_sel, col_elim = st.columns([1, 10, 1.5, 2])
            
            if i in st.session_state.eliminated_opts:
                col_num.markdown(f'<span class="option-number eliminated">{circle_numbers[i]}</span>', unsafe_allow_html=True)
            else:
                col_num.markdown(f'<span class="option-number">{circle_numbers[i]}</span>', unsafe_allow_html=True)

            text_style = "eliminated" if i in st.session_state.eliminated_opts else ""
            if st.session_state.selected_opt == i:
                col_text.markdown(f'<div class="{text_style}" style="font-weight: bold; color: #1971c2;">{opt_text}</div>', unsafe_allow_html=True)
            else:
                col_text.markdown(f'<div class="{text_style}">{opt_text}</div>', unsafe_allow_html=True)

            btn_label = "●" if st.session_state.selected_opt == i else "○"
            if col_sel.button(btn_label, key=f"sel_{idx}_{i}"):
                st.session_state.selected_opt = i
                st.rerun()

            elim_label = "해제" if i in st.session_state.eliminated_opts else "오답"
            btn_type = "secondary" if i in st.session_state.eliminated_opts else "primary"
            if col_elim.button(elim_label, key=f"elim_{idx}_{i}", type=btn_type):
                if i in st.session_state.eliminated_opts: st.session_state.eliminated_opts.remove(i)
                else: st.session_state.eliminated_opts.add(i)
                st.rerun()
        st.markdown('</div>', unsafe_allow_html=True)
        st.write("---")

        if st.button("🚀 정답 확인", type="primary", use_container_width=True, disabled=st.session_state.show_explanation):
            if st.session_state.selected_opt is None: st.warning("답을 선택해주세요!")
            else:
                st.session_state.show_explanation = True
                if st.session_state.selected_opt == card['correct_index']:
                    st.balloons(); st.success("✅ 정답입니다!"); update_card_schedule(idx, True)
                else:
                    st.error(f"❌ 오답입니다. 정답은 {circle_numbers[card['correct_index']]}번 입니다."); update_card_schedule(idx, False)
                st.rerun()

        if st.session_state.show_explanation:
            with st.expander("💡 해설 보기", expanded=True): st.info(card['explanation'])
            if st.button("➡️ 다음 문제 풀기", type="primary", use_container_width=True):
                st.session_state.show_explanation = False; st.rerun()

# ==========================================
# [탭 3] 문제 관리
# ==========================================
with tab3:
    st.header("🗂️ 문제 리스트")
    cards = load_cards()
    if not cards: st.write("저장된 문제가 없습니다.")
    else:
        circle_numbers = ["①", "②", "③", "④", "⑤"]
        selected_for_delete = []

        for i, card in enumerate(cards):
            col_exp, col_chk = st.columns([20, 1])
            with col_exp:
                with st.expander(f"#{i+1}. {card['question'][:50]}..."):
                    st.markdown(f'<div class="question-box">**Q.** {card["question"]}</div>', unsafe_allow_html=True)
                    st.markdown('<div class="options-box">', unsafe_allow_html=True)
                    for opt_i, opt_text in enumerate(card['options']):
                        if opt_i == card['correct_index']:
                            st.markdown(f'<div class="option-item" style="background-color: #e7f5ff;"><span class="option-number" style="color: #1971c2;">{circle_numbers[opt_i]}</span><span style="color: #1971c2; font-weight: bold;">{opt_text}</span></div>', unsafe_allow_html=True)
                        else:
                            st.markdown(f'<div class="option-item"><span class="option-number">{circle_numbers[opt_i]}</span><span>{opt_text}</span></div>', unsafe_allow_html=True)
                    st.markdown('</div>', unsafe_allow_html=True)
                    st.caption(f"💡 해설: {card['explanation']}")
            with col_chk:
                if st.checkbox("", key=f"chk_{i}", label_visibility="collapsed"):
                    selected_for_delete.append(i)

        st.divider()
        col_btn1, col_btn2 = st.columns([1, 1])
        with col_btn1:
            if st.button(f"🗑️ 선택 삭제 ({len(selected_for_delete)}개)", type="primary", use_container_width=True, disabled=len(selected_for_delete) == 0):
                remaining = [c for i, c in enumerate(cards) if i not in selected_for_delete]
                save_all_cards(remaining)
                st.rerun()
        with col_btn2:
            if st.button("🗑️ 전체 삭제", type="secondary", use_container_width=True):
                save_all_cards([])
                st.rerun()

# ==========================================
# [탭 4] 정리본 형성
# ==========================================
with tab4:
    lecture_content = ""
    jokbo_content = ""
    col_upload1, col_upload2 = st.columns(2)

    with col_upload1:
        uploaded_summaries = st.file_uploader("강의자료 업로드", type=['pdf', 'pptx'], key="summary_uploader", accept_multiple_files=True, label_visibility="collapsed")
        if uploaded_summaries:
            all_texts = []
            for f in uploaded_summaries:
                if f.name.endswith('.pdf'):
                    try:
                        reader = PdfReader(f)
                        text = "\n".join([p.extract_text() or "" for p in reader.pages])
                        if text.strip(): all_texts.append(text)
                    except: pass
                elif f.name.endswith('.pptx'):
                    try:
                        prs = PptxPresentation(f)
                        txt = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if shape.has_text_frame: txt.append(shape.text_frame.text)
                        all_texts.append("\n".join(txt))
                    except: pass
            lecture_content = "\n\n".join(all_texts)
            if lecture_content: st.success(f"강의자료 읽기 성공! ({len(lecture_content)}자)")


    with col_upload2:
        uploaded_jokbo = st.file_uploader("족보 업로드", type=['pdf', 'docx'], key="jokbo_uploader", label_visibility="collapsed")
        if uploaded_jokbo:
            if uploaded_jokbo.name.endswith('.pdf'):
                try:
                    reader = PdfReader(uploaded_jokbo)
                    jokbo_content = "\n".join([p.extract_text() or "" for p in reader.pages])
                except: pass
            elif uploaded_jokbo.name.endswith('.docx'):
                try:
                    doc = docx.Document(uploaded_jokbo)
                    jokbo_content = "\n".join([p.text for p in doc.paragraphs])
                except: pass
            if jokbo_content: st.success(f"족보 읽기 성공! ({len(jokbo_content)}자)")

    st.divider()

    if st.button("📋 통합 표 정리본 생성", type="primary", use_container_width=True, disabled=not bool(lecture_content)):
        with st.spinner("AI가 강의와 족보를 분석하여 표를 만들고 있습니다... (약 20초 소요)"):
            try:
                prompt = f"""
                당신은 의대 학습 정리 전문가입니다.
                강의자료를 메인 주제(질환 등)별로 나누고, 표 형태로 정리하세요.

                [구조 요구사항]
                1. 기본적으로 '소주제' - '내용'의 2단 구성을 따릅니다.
                2. 단, 소주제 내부에서 또다시 분류가 필요한 경우(예: 진단 내의 혈액검사/영상검사 등)에는 '세부 분류'를 추가하여 3단으로 구성하세요.
                
                [서식 규칙]
                1. 내용(value)은 긴 줄글로 쓰지 말고, 반드시 '1. ', '2. ' 번호를 붙여 개조식으로 작성하세요.
                2. 각 번호 항목이 끝날 때마다 반드시 줄바꿈을 하세요.
                
                [색상 태그 규칙]
                - 족보 정답 선지 내용: <yellow>내용</yellow>
                - 족보 오답 선지(강의 관련): <blue>내용</blue>
                - 족보 오답 선지(강의 무관): <gray>내용</gray>
                
                [입력 자료]
                강의: {lecture_content[:30000]}
                족보: {jokbo_content[:20000]}

                [출력 형식 - JSON 배열]
                반드시 아래 구조를 지키세요. 'sub_key'는 하위 분류가 있을 때만 작성하고, 없으면 null 또는 빈 문자열로 두세요.
                [
                  {{
                    "main_topic": "메인 주제명 (예: 급성 A형 간염)",
                    "sub_sections": [
                      {{ "key": "개요", "sub_key": "", "value": "1. 정의: ...\\n2. 역학: ..." }},
                      {{ "key": "진단", "sub_key": "혈액검사", "value": "1. IgM anti-HAV <yellow>양성</yellow>...\\n2. LFT 상승..." }},
                      {{ "key": "진단", "sub_key": "영상검사", "value": "1. 초음파: 간비대 소견..." }}
                    ]
                  }},
                  ...
                ]
                """
                # [핵심 수정] 무조건 JSON 형식으로만 강제 출력하도록 설정 추가
                response = client.models.generate_content(
                    model=MODEL, 
                    contents=prompt,
                    config={"response_mime_type": "application/json"}
                )
                
                st.session_state['summary_data'] = json.loads(response.text)
                st.rerun()
                
            except json.JSONDecodeError:
                st.error("AI가 올바른 형식(JSON)으로 표를 만들지 못했습니다. 다시 시도해 주세요.")
                with st.expander("AI 응답 원본 확인 (디버깅용)"):
                    st.write(response.text if 'response' in locals() else "응답 없음")
            except Exception as e: 
                st.error(f"오류: {e}")

    # ── 워드 다운로드 ──
    if st.session_state['summary_data']:
        st.success("✅ 정리본 생성이 완료되었습니다! 아래 버튼을 눌러 다운로드하세요.")
        
        try:
            doc_out = DocxDocument()
            
            # [제목]
            title = doc_out.add_heading('의대 강의/족보 통합 정리본', level=0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run_title = title.runs[0]
            set_font_style(run_title, font_size=16, is_bold=True)
            
            # [범례]
            legend = doc_out.add_paragraph()
            legend.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            run_y = legend.add_run('■ 정답  ')
            set_font_style(run_y, font_size=9)
            run_y.font.highlight_color = WD_COLOR_INDEX.YELLOW
            
            run_b = legend.add_run('■ 관련 오답  ')
            set_font_style(run_b, font_size=9)
            run_b.font.color.rgb = RGBColor(0x19, 0x71, 0xC2)
            
            run_g = legend.add_run('■ 무관 오답')
            set_font_style(run_g, font_size=9)
            run_g.font.color.rgb = RGBColor(0xAD, 0xB5, 0xBD)
            
            doc_out.add_paragraph() 

            for item in st.session_state['summary_data']:
                main_topic = item.get('main_topic', '')
                sub_sections = item.get('sub_sections', [])
                
                if not sub_sections: continue

                # 3열 테이블 생성 (소주제 / 세부분류 / 내용)
                table = doc_out.add_table(rows=0, cols=3)
                table.style = 'Table Grid' 
                
                # 메인 주제 행 (3칸 병합)
                row_main = table.add_row()
                cell_main = row_main.cells[0]
                cell_main.merge(row_main.cells[1])
                cell_main.merge(row_main.cells[2])
                cell_main.text = main_topic
                
                set_cell_background(cell_main, "495057") 
                p_main = cell_main.paragraphs[0]
                p_main.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run_main = p_main.runs[0]
                run_main.font.color.rgb = RGBColor(255, 255, 255)
                set_font_style(run_main, font_size=10, is_bold=True)

                last_key = None
                key_cell_anchor = None

                for sub in sub_sections:
                    key = sub.get('key', '')
                    sub_key = sub.get('sub_key', '')
                    content = sub.get('value', '')
                    
                    row = table.add_row()
                    row.height_rule = WD_ROW_HEIGHT_RULE.AT_LEAST
                    row.height = Cm(1.5) 
                    
                    # ── 1열: 소주제 (셀 병합 + 폰트 9pt) ──
                    cell_key = row.cells[0]
                    
                    if key == last_key and key_cell_anchor is not None:
                        key_cell_anchor.merge(cell_key)
                    else:
                        cell_key.text = key
                        cell_key.width = Cm(2.5) 
                        set_cell_background(cell_key, "E9ECEF")
                        
                        p_k = cell_key.paragraphs[0]
                        p_k.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        set_font_style(p_k.runs[0], font_size=9, is_bold=True)
                        cell_key.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
                        
                        key_cell_anchor = cell_key
                        last_key = key
                    
                    # ── 2열 & 3열 처리 ──
                    if sub_key and sub_key.strip():
                        cell_sub = row.cells[1]
                        cell_sub.text = sub_key
                        cell_sub.width = Cm(2.5)
                        set_cell_background(cell_sub, "F8F9FA")
                        
                        p_sub = cell_sub.paragraphs[0]
                        p_sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        set_font_style(p_sub.runs[0], font_size=9, is_bold=True)
                        cell_sub.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
                        
                        cell_val = row.cells[2]
                    else:
                        cell_sub = row.cells[1]
                        cell_sub.merge(row.cells[2])
                        cell_val = row.cells[1]

                    # ── 내용 채우기 ──
                    cell_val.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
                    p = cell_val.paragraphs[0]
                    p.paragraph_format.line_spacing = 1.0 
                    p.paragraph_format.space_before = Pt(6)
                    p.paragraph_format.space_after = Pt(6)
                    
                    parts = re.split(r'(<(?:yellow|blue|gray)>.*?</(?:yellow|blue|gray)>)', content)
                    for part in parts:
                        if not part: continue
                        tag_match = re.match(r'<(yellow|blue|gray)>(.*?)</\1>', part)
                        if tag_match:
                            tag_type = tag_match.group(1)
                            text_body = tag_match.group(2)
                            run = p.add_run(text_body)
                            
                            if tag_type == 'yellow':
                                run.font.highlight_color = WD_COLOR_INDEX.YELLOW
                                set_font_style(run, font_size=9, is_bold=False)
                            elif tag_type == 'blue':
                                run.font.color.rgb = RGBColor(0x19, 0x71, 0xC2)
                                set_font_style(run, font_size=9, is_bold=True)
                            elif tag_type == 'gray':
                                run.font.color.rgb = RGBColor(0xAD, 0xB5, 0xBD)
                                set_font_style(run, font_size=9, is_bold=False)
                        else:
                            run = p.add_run(part)
                            set_font_style(run, font_size=9, is_bold=False)

                doc_out.add_paragraph() 

            bio = BytesIO()
            doc_out.save(bio)
            bio.seek(0)
            
            st.download_button("💾 표 정리본 다운로드 (Word)", data=bio, file_name="통합_표_정리본.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)

        except Exception as e:
            st.error(f"워드 생성 오류: {e}")